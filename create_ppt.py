from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Privacy-Preserving Federated Learning for Disease Prediction"
    subtitle.text = "Team: [Your Team Name]\nHackWithUttarPradesh 2025 | IEDC Chandigarh University\nDate: [Current Date]"

    # Slide 2: Agenda
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Agenda'
    tf = body_shape.text_frame
    tf.text = 'Problem Statement'

    p = tf.add_paragraph()
    p.text = 'Solution Overview'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Technical Implementation'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Business Model'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Future Roadmap'
    p.level = 0

    # Slide 3: Problem Statement
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Problem Statement'
    tf = body_shape.text_frame
    tf.text = 'Healthcare Data Challenges'

    p = tf.add_paragraph()
    p.text = 'Data Privacy Concerns: Patient data is highly sensitive and regulated (HIPAA, GDPR)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Data Silos: Hospitals cannot share patient data due to privacy laws'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Limited Training Data: Individual hospitals have insufficient data for robust AI models'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Lack of Transparency: Black-box AI models reduce trust in clinical decisions'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Current Issues'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Centralized ML: Requires data aggregation, violating privacy'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Bias in Models: Limited datasets lead to biased predictions'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'No Explainability: Clinicians cannot understand AI decisions'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Security Risks: Data breaches during transmission'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Impact: Reduced adoption of AI in healthcare despite potential benefits'
    p.level = 0

    # Slide 4: Solution Overview
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Solution Overview'
    tf = body_shape.text_frame
    tf.text = 'Federated Learning Approach'

    p = tf.add_paragraph()
    p.text = 'Distributed Training: Train AI models across multiple hospitals without sharing raw data'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Privacy Preservation: Use differential privacy to protect individual patient information'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Explainable AI: Provide SHAP and LIME explanations for clinical interpretability'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Real-time Collaboration: Hospitals collaborate on model improvement while maintaining data sovereignty'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Key Benefits'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Privacy-First: No raw patient data leaves hospital premises'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Scalable: Leverages data from multiple institutions'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Trustworthy: Transparent AI decisions with explanations'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Regulatory Compliant: Meets HIPAA and GDPR requirements'
    p.level = 1

    # Slide 5: Solution Architecture (Text-based diagram)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Solution Architecture'
    tf = body_shape.text_frame
    tf.text = 'Federated Learning System Architecture'

    p = tf.add_paragraph()
    p.text = 'Hospital 1 (Client) ────┐'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Hospital 2 (Client) ────┼─── FL Server (Aggregator) ──── Global Model + XAI'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Hospital 3 (Client) ────┘'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Local Data • DP Training • FedAvg • Privacy Check • SHAP + LIME'
    p.level = 1

    # Slide 6: Technical Implementation
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Technical Implementation'
    tf = body_shape.text_frame
    tf.text = 'Core Technologies'

    p = tf.add_paragraph()
    p.text = 'Federated Learning: Flower framework for client-server architecture'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Privacy: Opacus library for differential privacy (ε=1.0, δ=1e-5)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Explainability: SHAP for global explanations, LIME for local explanations'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Model: PyTorch-based MLP classifier with configurable architecture'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Data: UCI Heart Disease and Diabetes datasets'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Key Features'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Multi-Client Simulation: 3+ simulated hospitals with heterogeneous data'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Secure Aggregation: Encrypted model weight updates'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Gradient Clipping: L2 norm clipping to prevent privacy leakage'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Noise Injection: Gaussian noise calibrated to privacy budget'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Performance Metrics'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Accuracy: 80-85% on test datasets'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'AUROC: 85-90% for disease prediction'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Privacy Budget: Tracks cumulative privacy loss'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Explanation Stability: Consistent interpretations across FL rounds'
    p.level = 1

    # Slide 7: Clinical Applications
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Clinical Applications'
    tf = body_shape.text_frame
    tf.text = 'Disease Prediction Models'

    p = tf.add_paragraph()
    p.text = 'Cardiovascular Risk: Early detection using clinical markers'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Diabetes Prediction: Risk assessment from patient history'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Dermatological Analysis: Skin condition diagnosis with severity scoring'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Clinical Features'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'PASI Score Calculation: Psoriasis severity assessment'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'SCORAD Index: Atopic dermatitis evaluation'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Evidence-Based Treatments: Protocol-driven therapy recommendations'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Follow-up Planning: Structured monitoring schedules'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Real-World Impact'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Early Intervention: Proactive disease management'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Resource Optimization: Efficient healthcare resource allocation'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Clinical Decision Support: AI-assisted diagnosis and treatment'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Patient Outcomes: Improved health outcomes through timely interventions'
    p.level = 1

    # Slide 8: Technical Stack
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Technical Stack'
    tf = body_shape.text_frame
    tf.text = 'Backend'

    p = tf.add_paragraph()
    p.text = 'Python 3.9+'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'PyTorch: Deep learning framework'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Flower: Federated learning platform'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Opacus: Differential privacy library'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Frontend'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Streamlit: Interactive web dashboard'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Plotly: Data visualization'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Matplotlib/Seaborn: Statistical plots'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Data & Privacy'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Pandas/NumPy: Data processing'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Scikit-learn: Traditional ML algorithms'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'SHAP/LIME: Explainability frameworks'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'PyCryptodome: Encryption utilities'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Deployment'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Docker: Containerization'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'AWS/Azure/GCP: Cloud deployment options'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Streamlit Cloud: Free hosting platform'
    p.level = 1

    # Slide 9: Business Model
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Business Model'
    tf = body_shape.text_frame
    tf.text = 'Revenue Streams'

    p = tf.add_paragraph()
    p.text = '1. SaaS Subscription Model'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Tiered Pricing: Basic ($99/month), Professional ($299/month), Enterprise ($999/month)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Features: Model training, privacy monitoring, basic support'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '2. Deployment Services'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Implementation Fee: $5,000-15,000 per hospital'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Customization: Dataset integration, model fine-tuning'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Training: Staff training and onboarding'
    p.level = 2

    p = tf.add_paragraph()
    p.text = '3. Premium Features'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Advanced Analytics: $199/month'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Custom Models: $499/month'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Priority Support: $99/month'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Target Market'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Primary: Large hospital networks (500+ beds)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Secondary: Regional healthcare systems'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Tertiary: Research institutions and pharma companies'
    p.level = 1

    # Slide 10: Market Opportunity
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Market Opportunity'
    tf = body_shape.text_frame
    tf.text = 'Market Size'

    p = tf.add_paragraph()
    p.text = 'Global Healthcare AI Market: $45B by 2026 (CAGR 49%)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Federated Learning in Healthcare: $2.5B opportunity'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Privacy-Preserving AI: Growing regulatory demand'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Competitive Advantage'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Privacy-First Approach: Unique selling proposition'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Explainable AI: Regulatory compliance advantage'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Open-Source Foundation: Lower development costs'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Clinical Validation: Medical-grade reliability'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Go-to-Market Strategy'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Phase 1: Pilot with 3-5 hospitals (6 months)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Phase 2: Regional expansion (12-18 months)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Phase 3: National healthcare network (24+ months)'
    p.level = 1

    # Slide 11: Financial Projections
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Financial Projections'
    tf = body_shape.text_frame
    tf.text = 'Year 1 Projections'

    p = tf.add_paragraph()
    p.text = 'Revenue: $500K (pilot programs + subscriptions)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Customers: 10 hospitals'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Cost: $300K (development + operations)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Profit: $200K'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Year 3 Projections'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Revenue: $5M (expanded customer base)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Customers: 200+ hospitals'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Cost: $1.5M (scaling infrastructure)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Profit: $3.5M'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Funding Requirements'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Seed Round: $1M for product development'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Series A: $5M for market expansion'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Break-even: 18 months post-launch'
    p.level = 1

    # Slide 12: Risk Analysis & Mitigation
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Risk Analysis & Mitigation'
    tf = body_shape.text_frame
    tf.text = 'Technical Risks'

    p = tf.add_paragraph()
    p.text = 'Privacy Vulnerabilities: Regular security audits, third-party validation'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Model Performance: Continuous monitoring, retraining pipelines'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Scalability Issues: Cloud-native architecture, load testing'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Regulatory Risks'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Compliance Changes: Legal team monitoring, adaptive compliance'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Data Privacy Laws: Privacy-by-design approach, GDPR/HIPAA compliance'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Market Risks'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Adoption Resistance: Clinical validation studies, physician training'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Competition: First-mover advantage, proprietary algorithms'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Mitigation Strategy'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Pilot Testing: Extensive validation with medical partners'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Insurance: Cyber liability and professional liability coverage'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Backup Plans: Centralized fallback options, data anonymization'
    p.level = 1

    # Slide 13: Future Roadmap
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Future Roadmap'
    tf = body_shape.text_frame
    tf.text = 'Phase 1: MVP (Current)'

    p = tf.add_paragraph()
    p.text = '✅ Federated learning pipeline'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '✅ Differential privacy implementation'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '✅ Basic explainability features'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '✅ Streamlit dashboard'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Phase 2: Clinical Validation (3-6 months)'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '🔄 IRB-approved clinical trials'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '🔄 Physician validation studies'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '🔄 Regulatory compliance (FDA/CE)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '🔄 Real hospital integration'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Phase 3: Advanced Features (6-12 months)'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '🔄 Multi-modal data integration (images, text)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '🔄 Advanced privacy mechanisms'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '🔄 Real-time federated inference'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '🔄 Mobile application'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Phase 4: Enterprise Scale (12-24 months)'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '🔄 Global healthcare network'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '🔄 AI model marketplace'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '🔄 Research collaboration platform'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '🔄 International expansion'
    p.level = 1

    # Slide 14: Team & Expertise
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Team & Expertise'
    tf = body_shape.text_frame
    tf.text = 'Core Team'

    p = tf.add_paragraph()
    p.text = 'AI/ML Engineer: Federated learning and privacy expertise'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Clinical Advisor: Medical doctor with AI experience'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Security Expert: Privacy and compliance specialist'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Product Manager: Healthcare technology background'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Advisors'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Healthcare Policy: Regulatory compliance guidance'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Data Privacy: GDPR/HIPAA expertise'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Clinical Research: Medical validation support'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Partnerships'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Academic Institutions: Research collaboration'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Healthcare Providers: Pilot program partners'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Technology Vendors: Infrastructure support'
    p.level = 1

    # Slide 15: Call to Action
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Call to Action'
    tf = body_shape.text_frame
    tf.text = 'Next Steps'

    p = tf.add_paragraph()
    p.text = '1. Pilot Program: Partner with 3 hospitals for validation'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '2. Clinical Trials: IRB-approved studies for regulatory approval'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '3. Funding Round: Secure seed investment for development'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '4. Team Expansion: Hire clinical and technical experts'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Investment Opportunity'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Problem: Privacy-preserving AI in healthcare'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Solution: Production-ready federated learning platform'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Market: $45B healthcare AI market'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Traction: Validated technology, clinical applications'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Join us in revolutionizing healthcare with privacy-preserving AI!'
    p.level = 0

    # Slide 16: Contact Information
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Contact Information'
    tf = body_shape.text_frame
    tf.text = 'Project Lead: [Your Name]'

    p = tf.add_paragraph()
    p.text = 'Email: [your.email@example.com]'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'LinkedIn: [linkedin.com/in/yourprofile]'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'GitHub: [github.com/yourusername]'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Project Repository: [github.com/yourusername/ai-disease-prediction]'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Demo Link: [streamlit deployment URL]'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Thank You!'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Questions & Discussion'
    p.level = 0

    # Save the presentation
    prs.save('AI_Disease_Prediction_Presentation.pptx')
    print("Presentation created successfully: AI_Disease_Prediction_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
